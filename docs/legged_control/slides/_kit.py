# -*- coding: utf-8 -*-
"""Shared widescreen slide helpers for Japanese teaching decks."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1B, 0x3A, 0x4B)
TEAL = RGBColor(0x2A, 0x6F, 0x6F)
TERR = RGBColor(0xC4, 0x5A, 0x20)
GOLD = RGBColor(0xB4, 0x8A, 0x28)
INK = RGBColor(0x1E, 0x1E, 0x24)
MUTED = RGBColor(0x5A, 0x60, 0x6E)
PAPER = RGBColor(0xF7, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD4, 0xCF, 0xC6)
SOFT = RGBColor(0xEE, 0xEA, 0xE2)

FONT = "Noto Sans CJK JP"
W = Inches(13.333)
H = Inches(7.5)

EQ_DIR = Path(__file__).resolve().parent / "_eqcache"


def _set_run_font(run, size_pt, bold=False, color=INK, name=FONT, italic=False):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        node = rPr.find(qn(tag))
        if node is None:
            node = etree.SubElement(rPr, qn(tag))
        node.set("typeface", name)
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", name)


def add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold=bold, color=color)
    return box


def add_para(tf, text, size=16, bold=False, color=INK, space_before=6, space_after=2, level=0, bullet=False):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].text) else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        pPr = p._p.get_or_add_pPr()
        bu = pPr.find(qn("a:buFont"))
        if bu is None:
            etree.SubElement(pPr, qn("a:buChar")).set("char", "•")
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold=bold, color=color)
    return p


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(sh, color)
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def rounded(slide, left, top, width, height, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(sh, color)
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


class Deck:
    def __init__(self, footer: str):
        self.prs = Presentation()
        self.prs.slide_width = W
        self.prs.slide_height = H
        self.footer = footer
        self._count = 0
        EQ_DIR.mkdir(parents=True, exist_ok=True)

    def blank(self):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        rect(slide, 0, 0, W, H, PAPER)
        self._count += 1
        return slide

    def _footer(self, slide):
        rect(slide, 0, Inches(7.22), W, Inches(0.28), NAVY)
        add_text_box(slide, Inches(0.4), Inches(7.22), Inches(10.5), Inches(0.28), self.footer, size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.28), str(self._count), size=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    def _bar(self, slide, kicker: str):
        rect(slide, 0, 0, W, Inches(0.92), NAVY)
        rect(slide, 0, Inches(0.92), W, Inches(0.06), TEAL)
        add_text_box(slide, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.28), kicker, size=11, color=RGBColor(0xB8, 0xD4, 0xD4), bold=True)
        return slide

    def title_slide(self, title: str, subtitle: str, meta: str):
        slide = self.blank()
        rect(slide, 0, 0, Inches(0.22), H, TEAL)
        rect(slide, 0, 0, W, Inches(0.18), NAVY)
        add_text_box(slide, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4), "docs/legged_control  /  学習ノートのスライド化", size=16, color=TEAL, bold=True)
        add_text_box(slide, Inches(0.7), Inches(2.15), Inches(12), Inches(1.6), title, size=36, color=NAVY, bold=True)
        add_text_box(slide, Inches(0.7), Inches(3.9), Inches(11.5), Inches(1.2), subtitle, size=18, color=INK)
        add_text_box(slide, Inches(0.7), Inches(6.4), Inches(11.5), Inches(0.5), meta, size=13, color=MUTED)
        return slide

    def section_slide(self, code: str, title: str, blurb: str):
        slide = self.blank()
        rect(slide, 0, 0, W, H, NAVY)
        rect(slide, 0, 0, Inches(0.22), H, TEAL)
        add_text_box(slide, Inches(0.7), Inches(2.1), Inches(12), Inches(0.4), code, size=16, color=RGBColor(0xB8, 0xD4, 0xD4), bold=True)
        add_text_box(slide, Inches(0.7), Inches(2.55), Inches(12), Inches(1.3), title, size=34, color=WHITE, bold=True)
        add_text_box(slide, Inches(0.7), Inches(4.1), Inches(11.5), Inches(1.4), blurb, size=18, color=RGBColor(0xD7, 0xE0, 0xE0))
        add_text_box(slide, Inches(12.4), Inches(6.9), Inches(0.7), Inches(0.3), str(self._count), size=11, color=WHITE, align=PP_ALIGN.RIGHT)
        return slide

    def content(self, kicker: str, title: str, takeaway: str | None = None):
        slide = self.blank()
        self._bar(slide, kicker)
        add_text_box(slide, Inches(0.4), Inches(0.36), Inches(12.5), Inches(0.5), title, size=24, color=WHITE, bold=True)
        if takeaway:
            rounded(slide, Inches(0.4), Inches(1.12), Inches(12.5), Inches(0.55), SOFT, line=LINE)
            add_text_box(slide, Inches(0.55), Inches(1.16), Inches(12.2), Inches(0.48), "要点  " + takeaway, size=14, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        self._footer(slide)
        return slide

    def fn(self, slide, text: str, top=6.40):
        """Bottom ※ notes. Footer bar starts at 7.22."""
        rounded(slide, Inches(0.35), Inches(top), Inches(12.63), Inches(0.76), SOFT, line=LINE)
        add_text_box(slide, Inches(0.45), Inches(top + 0.03), Inches(12.43), Inches(0.70), text, size=11, color=MUTED)
        return slide

    def pre(self, slide, text, left=0.4, top=1.2, width=12.5, height=5.6, size=11):
        """Monospace block for the ASCII diagrams that are already in the notes."""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = False
        first = True
        for line in text.splitlines():
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = line if line else " "
            _set_run_font(run, size, color=INK, name="DejaVu Sans Mono")
        return box

    def bullets(self, slide, items, left=0.5, top=1.85, width=12.3, height=5.0, size=17):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            if isinstance(item, tuple):
                text, opts = item[0], item[1] if len(item) > 1 else {}
            else:
                text, opts = item, {}
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = opts.get("level", 0)
            p.space_before = Pt(opts.get("sb", 7))
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = ("•  " if opts.get("level", 0) == 0 else "–  ") + text
            _set_run_font(run, opts.get("size", size), bold=opts.get("bold", False), color=opts.get("color", INK))
        return box

    def badge_row(self, slide, badges, top=1.85):
        x = Inches(0.45)
        for label, color in badges:
            w = Inches(0.28 * (len(label) + 2))
            rounded(slide, x, Inches(top), w, Inches(0.34), color)
            add_text_box(slide, x, Inches(top), w, Inches(0.34), label, size=11, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
            x += w + Inches(0.12)

    def cards(self, slide, cards, top=1.9, height=4.7):
        n = len(cards)
        gap = 0.18
        left0 = 0.4
        usable = 12.5 - gap * (n - 1)
        cw = usable / n
        for i, (head, lines, accent) in enumerate(cards):
            x = Inches(left0 + i * (cw + gap))
            rounded(slide, x, Inches(top), Inches(cw), Inches(height), WHITE, line=LINE)
            rect(slide, x, Inches(top), Inches(cw), Inches(0.08), accent)
            add_text_box(slide, x + Inches(0.16), Inches(top + 0.18), Inches(cw - 0.3), Inches(0.7), head, size=15, color=NAVY, bold=True)
            box = slide.shapes.add_textbox(x + Inches(0.16), Inches(top + 0.9), Inches(cw - 0.3), Inches(height - 1.1))
            tf = box.text_frame
            tf.word_wrap = True
            for j, line in enumerate(lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_before = Pt(5)
                run = p.add_run()
                run.text = line
                _set_run_font(run, 13, color=INK)

    def table(self, slide, rows, left, top, width, height, col_w=None, font=12):
        n_r, n_c = len(rows), len(rows[0])
        tbl = slide.shapes.add_table(n_r, n_c, Inches(left), Inches(top), Inches(width), Inches(height)).table
        if col_w:
            for i, w in enumerate(col_w):
                tbl.columns[i].width = Inches(w)
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = ""
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                header = r == 0
                _set_run_font(run, font, bold=header, color=WHITE if header else INK)
                fill = NAVY if header else (SOFT if r % 2 else WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
        return tbl

    def eq(self, latex: str, name: str, fontsize: float = 22) -> Path:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        path = EQ_DIR / f"{name}.png"
        fig = plt.figure(figsize=(11, 1.15))
        fig.patch.set_facecolor("white")
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis("off")
        ax.text(0.5, 0.5, f"${latex}$", ha="center", va="center", fontsize=fontsize)
        fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return path

    def add_eq(self, slide, latex, name, left=0.6, top=3.2, width=12.1, fontsize=22):
        path = self.eq(latex, name, fontsize)
        rounded(slide, Inches(left), Inches(top), Inches(width), Inches(1.25), WHITE, line=LINE)
        slide.shapes.add_picture(str(path), Inches(left + 0.15), Inches(top + 0.12), width=Inches(width - 0.3))

    def bridge(self, slide, prev: str, now: str, nxt: str, top=1.12):
        """Show how this page sits between the previous exit and the next entrance."""
        items = [("前の出口", prev, TEAL), ("このページ", now, NAVY), ("次の入口", nxt, TERR)]
        gap = 0.14
        cw = (12.5 - 2 * gap) / 3
        for i, (lab, text, col) in enumerate(items):
            x = Inches(0.4 + i * (cw + gap))
            rounded(slide, x, Inches(top), Inches(cw), Inches(0.95), WHITE, line=LINE)
            rect(slide, x, Inches(top), Inches(cw), Inches(0.26), col)
            add_text_box(slide, x, Inches(top), Inches(cw), Inches(0.26), lab, size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + Inches(0.08), Inches(top + 0.3), Inches(cw - 0.16), Inches(0.6), text, size=12, color=INK)

    def io_sentence(self, slide, text: str, top=2.15):
        rounded(slide, Inches(0.4), Inches(top), Inches(12.5), Inches(0.7), SOFT, line=LINE)
        add_text_box(slide, Inches(0.55), Inches(top), Inches(12.2), Inches(0.7), text, size=16, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    def eq_lesson(self, kicker: str, title: str, latex: str, name: str, background: str, intent: str, variables: list, fontsize: float = 20, notes: str | None = None):
        """One equation per slide: 背景, 意図, 式, 変数表, 下部※."""
        slide = self.content(kicker, title)
        rounded(slide, Inches(0.4), Inches(1.12), Inches(6.2), Inches(1.15), SOFT, line=LINE)
        add_text_box(slide, Inches(0.52), Inches(1.16), Inches(5.96), Inches(0.28), "背景", size=12, color=TEAL, bold=True)
        add_text_box(slide, Inches(0.52), Inches(1.42), Inches(5.96), Inches(0.78), background, size=13, color=INK)
        rounded(slide, Inches(6.75), Inches(1.12), Inches(6.15), Inches(1.15), SOFT, line=LINE)
        add_text_box(slide, Inches(6.87), Inches(1.16), Inches(5.91), Inches(0.28), "意図", size=12, color=NAVY, bold=True)
        add_text_box(slide, Inches(6.87), Inches(1.42), Inches(5.91), Inches(0.78), intent, size=13, color=INK)
        self.add_eq(slide, latex, name, left=0.4, top=2.4, width=12.5, fontsize=fontsize)
        self.table(slide, variables, 0.4, 3.85, 12.5, min(2.40, 0.40 * len(variables)), font=13)
        if notes:
            self.fn(slide, notes)
        return slide

    def frame4(self, slide, bg, intent, arch, flow, top=1.2):
        """2x2 educational frame: 背景 / 意図 / アーキ / データの流れ."""
        cells = [
            ("背景  なぜこの問題があるか", bg, TEAL),
            ("意図  何を達成したいか", intent, NAVY),
            ("アーキ  どこに置き、何と繋ぐか", arch, TERR),
            ("データの流れ  何が入り、何が出るか", flow, GOLD),
        ]
        gap_x, gap_y = 0.18, 0.16
        cw, ch = (12.5 - gap_x) / 2, (5.65 - gap_y) / 2
        for i, (head, lines, accent) in enumerate(cells):
            r, c = divmod(i, 2)
            x = Inches(0.4 + c * (cw + gap_x))
            y = Inches(top + r * (ch + gap_y))
            rounded(slide, x, y, Inches(cw), Inches(ch), WHITE, line=LINE)
            rect(slide, x, y, Inches(0.1), Inches(ch), accent)
            add_text_box(slide, x + Inches(0.22), y + Inches(0.1), Inches(cw - 0.35), Inches(0.38), head, size=14, color=accent, bold=True)
            box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.5), Inches(cw - 0.38), Inches(ch - 0.62))
            tf = box.text_frame
            tf.word_wrap = True
            for j, line in enumerate(lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = "•  " + line
                _set_run_font(run, 13, color=INK)

    def node(self, slide, l, t, w, h, tag, title, sub="", accent=TEAL):
        """Labeled architecture box. Coordinates in inches."""
        rounded(slide, Inches(l), Inches(t), Inches(w), Inches(h), WHITE, line=LINE)
        rect(slide, Inches(l), Inches(t), Inches(0.09), Inches(h), accent)
        add_text_box(slide, Inches(l + 0.16), Inches(t + 0.04), Inches(w - 0.24), Inches(0.22), tag, size=10, color=accent, bold=True)
        add_text_box(slide, Inches(l + 0.16), Inches(t + 0.24), Inches(w - 0.24), Inches(0.36), title, size=13, color=NAVY, bold=True)
        if sub:
            add_text_box(slide, Inches(l + 0.16), Inches(t + 0.58), Inches(w - 0.24), Inches(h - 0.64), sub, size=11, color=INK)

    def arr(self, slide, x, y, text="→", w=0.36, h=0.28, color=TEAL):
        add_text_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), text, size=16, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def cap(self, slide, x, y, w, text, size=10, color=MUTED):
        add_text_box(slide, Inches(x), Inches(y), Inches(w), Inches(0.32), text, size=size, color=color)

    def hline(self, slide, x, y, w, color=TEAL):
        rect(slide, Inches(x), Inches(y), Inches(w), Inches(0.025), color)

    def vline(self, slide, x, y, h, color=TEAL):
        rect(slide, Inches(x), Inches(y), Inches(0.025), Inches(h), color)

    def vec_bar(self, slide, l, t, w, h, name, segs):
        """Horizontal stacked bar. segs = [(label, count, color), ...]."""
        add_text_box(slide, Inches(l), Inches(t), Inches(1.55), Inches(h), name, size=12, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        x = l + 1.65
        usable = w - 1.65
        total = sum(n for _, n, _ in segs) or 1
        for lab, n, col in segs:
            wi = usable * n / total
            rect(slide, Inches(x), Inches(t), Inches(wi - 0.04), Inches(h), col)
            add_text_box(
                slide,
                Inches(x),
                Inches(t),
                Inches(wi - 0.04),
                Inches(h),
                lab,
                size=11,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += wi

    def pipeline(self, slide, steps, top=2.05, box_h=1.55):
        """Horizontal architecture boxes. steps = [(tag, title, sub), ...]."""
        n = len(steps)
        gap = 0.12
        left0 = 0.35
        usable = 12.6 - gap * (n - 1)
        cw = usable / n
        for i, (tag, title, sub) in enumerate(steps):
            x = Inches(left0 + i * (cw + gap))
            rounded(slide, x, Inches(top), Inches(cw), Inches(box_h), WHITE, line=LINE)
            rect(slide, x, Inches(top), Inches(cw), Inches(0.32), TEAL if i < n - 1 else TERR)
            add_text_box(slide, x, Inches(top), Inches(cw), Inches(0.32), tag, size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + Inches(0.06), Inches(top + 0.38), Inches(cw - 0.12), Inches(0.55), title, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            add_text_box(slide, x + Inches(0.06), Inches(top + 0.92), Inches(cw - 0.12), Inches(0.55), sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)
            if i < n - 1:
                add_text_box(
                    slide,
                    x + Inches(cw - 0.08),
                    Inches(top + box_h / 2 - 0.16),
                    Inches(0.28),
                    Inches(0.32),
                    ">",
                    size=16,
                    color=TEAL,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )

    def pic(self, slide, path, left, top, width, height=None):
        """Place a PNG. height is optional; keep aspect if omitted."""
        rounded(slide, Inches(left - 0.04), Inches(top - 0.04), Inches(width + 0.08), Inches((height or 4.8) + 0.08), WHITE, line=LINE)
        kwargs = {"width": Inches(width)}
        if height is not None:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)

    def eq_with_fig(self, kicker, title, fig_path, latex, name, relation, notes=None, fontsize=16, fig_w=6.9, fig_h=5.05):
        """Figure left, equation right: 犬 / アーキ と式の対応。"""
        slide = self.content(kicker, title)
        self.pic(slide, fig_path, 0.38, 1.12, fig_w, fig_h)
        self.add_eq(slide, latex, name, left=7.45, top=1.18, width=5.50, fontsize=fontsize)
        rounded(slide, Inches(7.45), Inches(2.60), Inches(5.50), Inches(3.55), SOFT, line=LINE)
        add_text_box(slide, Inches(7.58), Inches(2.68), Inches(5.24), Inches(0.28), "式との関係", size=12, color=TEAL, bold=True)
        add_text_box(slide, Inches(7.58), Inches(2.98), Inches(5.24), Inches(3.05), relation, size=13, color=INK)
        if notes:
            self.fn(slide, notes)
        return slide

    def save(self, path: Path):
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path
